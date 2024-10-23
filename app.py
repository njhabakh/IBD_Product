import streamlit as st
import openai
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os
import requests
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
import sqlalchemy
from sqlalchemy import create_engine
from langchain.chat_models import AzureChatOpenAI

from dotenv import load_dotenv
load_dotenv(".env")

os.environ["AZURE_OPENAI_API_KEY"]     = os.environ.get("AZURE_OPENAI_API_KEY")
os.environ["AZURE_OPENAI_ENDPOINT"]    = os.environ.get("AZURE_OPENAI_ENDPOINT")
os.environ["AZURE_OPENAI_API_VERSION"] = os.environ.get("AZURE_OPENAI_API_VERSION")

def create_vector_database(txt_path):
    loader = TextLoader(txt_path)
    docs = loader.load()
    documents = RecursiveCharacterTextSplitter(
        chunk_size=1000, separators=["\n","\n\n"], chunk_overlap=200
    ).split_documents(docs)

    embeddings = AzureOpenAIEmbeddings(
        azure_deployment=os.environ.get("OPENAI_API_BASE"),
        openai_api_version=os.environ.get("OPENAI_API_VERSION"),
    )
    db = FAISS.from_documents(
        documents=documents,
        embedding=embeddings
    )
    db.save_local("./faiss-db")


def main():
    st.title("Investment Banking Analyst Toolkit")

    # Each tab runs independently
    tab1, tab2, tab3 = st.tabs(["Pitch Book Generator", "Q&A bot", "Valuation Models"])

    with tab1:
            try:
                pitch_book_generator()
            except Exception as e:
                st.error(f"Error in Pitch Book Generator: {e}")

    with tab2:
        try:
            chat_pod()
        except Exception as e:
            st.error(f"Error in Q&A bot: {e}")

    with tab3:
        try:
            valuation_models()
        except Exception as e:
            st.error(f"Error in Valuation Models: {e}")

def pitch_book_generator():
    st.header("Pitch Book Generator")
    st.write("Select the slides you want to include in your pitch book:")

    # Use columns to display controls and slides side by side
    col1, col2 = st.columns(2)

    with col1:
        slides = {
            "Upload Files and News": False,  # New slide type for data-driven content
            "Company Overview": False,
            "Market Analysis": False,
            "Financial Summary": False,
            "Competitor Analysis": False,
            "Investment Highlights": False,
            "Risks and Mitigations": False,            
        }

        for slide in slides.keys():
            slides[slide] = st.checkbox(slide)

        # Add custom slide option
        include_custom_slide = st.checkbox("Include Custom Slide")
        custom_slide_prompt = ""
        custom_slide_title = ""
        if include_custom_slide:
            custom_slide_title = st.text_input("Custom Slide Title", value="Custom Slide")
            custom_slide_prompt = st.text_area("Enter your custom prompt for the slide content")

        if slides["Upload Files and News"]:
            st.write("Upload documents and enter news links:")
            uploaded_files = st.file_uploader("Upload PDFs or text files", accept_multiple_files=True, type=['pdf', 'txt'])
            news_links = st.text_area("Enter news article URLs (one per line)")

            # Database connection details
            st.write("Enter SQL database connection details:")
            db_type = st.selectbox("Database Type", ["MySQL", "PostgreSQL", "SQLite"])
            db_host = st.text_input("Host", value="localhost")
            db_port = st.text_input("Port", value="3306" if db_type == "MySQL" else "5432")
            db_name = st.text_input("Database Name")
            db_user = st.text_input("Username")
            db_password = st.text_input("Password", type="password")
            db_query = st.text_area("Enter SQL Query to retrieve data")

        if st.button("Generate Pitch Book"):
            # Generate content for selected slides
            selected_slides = [slide for slide, selected in slides.items() if selected]

            if include_custom_slide and custom_slide_prompt.strip():
                selected_slides.append(custom_slide_title)

            if selected_slides:
                st.write("Generating the following slides:")
                prs = Presentation()
                # Use the second column to display slides
                with col2:
                    for slide in selected_slides:
                        st.write(f"- {slide}")
                        if slide == custom_slide_title:
                            content, charts = generate_slide_content(slide, custom_prompt=custom_slide_prompt)
                        elif slide == "Upload Files and News":
                            # Process uploaded files, news links, and database query
                            data_content, data_charts = generate_data_driven_slide(
                                uploaded_files, news_links, db_type, db_host, db_port, db_name, db_user, db_password, db_query
                            )
                            content = data_content
                            charts = data_charts
                        else:
                            content, charts = generate_slide_content(slide)
                        # Display content
                        st.subheader(slide)
                        st.write(content)
                        # Display charts if any
                        for fig in charts:
                            st.pyplot(fig)
                        # Add slide to PPT
                        slide_layout = prs.slide_layouts[5]  # Blank slide
                        ppt_slide = prs.slides.add_slide(slide_layout)
                        # Add title
                        txBox = ppt_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                        tf = txBox.text_frame
                        tf.text = slide
                        # Add content
                        left = Inches(1)
                        top = Inches(1.5)
                        width = Inches(8)
                        height = Inches(4)
                        txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = content
                        # Add charts to PPT
                        for fig in charts:
                            # Save figure temporarily
                            fig_file = f"{slide}.png"
                            fig.savefig(fig_file)
                            # Insert image into slide
                            left = Inches(1)
                            top = Inches(3)
                            pic = ppt_slide.shapes.add_picture(fig_file, left, top)
                            # Remove temp file
                            os.remove(fig_file)
                # Save PPT to a BytesIO object
                ppt_io = BytesIO()
                prs.save(ppt_io)
                ppt_io.seek(0)
                # Download button
                st.download_button(
                    label="Download Pitch Book",
                    data=ppt_io,
                    file_name="PitchBook.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.write("Please select at least one slide or provide a custom slide.")

def generate_slide_content(slide_name, custom_prompt=None):
    prompt = ""
    content = ""
    charts = []
    if custom_prompt:
        prompt = custom_prompt
    elif slide_name == "Company Overview":
        prompt = "Provide a detailed company overview for a pitch book presentation."
    elif slide_name == "Market Analysis":
        prompt = "Provide a market analysis including key trends, market size, and growth projections."
    elif slide_name == "Financial Summary":
        prompt = "Provide a financial summary including revenue, EBITDA, and net income over the past 5 years."
    elif slide_name == "Competitor Analysis":
        prompt = "Provide a competitor analysis including market share and key competitive advantages."
    elif slide_name == "Investment Highlights":
        prompt = "List the top investment highlights that make this company attractive."
    elif slide_name == "Risks and Mitigations":
        prompt = "List the key risks associated with the investment and potential mitigations."
    else:
        prompt = f"Provide information on {slide_name} for a pitch book presentation."


    llm = AzureChatOpenAI(
    azure_deployment=os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT"),
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    openai_api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
    verbose=False,
    temperature=0.3,
)

    from langchain.schema import (
    SystemMessage,
    HumanMessage,
    AIMessage
    )
    messages = [
        SystemMessage(content=f"You are an assistant helping with investment banking slides.  {prompt}" )                
    ]

    response = llm(
        messages= messages,
        max_tokens=100,
        temperature=0.2,
    )

    content = response.content

    # Generate charts if applicable
    if slide_name == "Market Analysis" and not custom_prompt:
        # Sample data for chart
        years = np.arange(2018, 2023)
        market_size = [100, 120, 140, 160, 180]
        fig, ax = plt.subplots()
        ax.plot(years, market_size, marker='o')
        ax.set_title('Market Size Over Years')
        ax.set_xlabel('Year')
        ax.set_ylabel('Market Size (in millions)')
        plt.tight_layout()
        charts.append(fig)
    elif slide_name == "Financial Summary" and not custom_prompt:
        years = np.arange(2018, 2023)
        revenue = [200, 220, 240, 260, 280]
        ebitda = [50, 55, 60, 65, 70]
        net_income = [30, 33, 36, 39, 42]
        fig, ax = plt.subplots()
        ax.plot(years, revenue, label='Revenue')
        ax.plot(years, ebitda, label='EBITDA')
        ax.plot(years, net_income, label='Net Income')
        ax.set_title('Financial Performance Over Years')
        ax.set_xlabel('Year')
        ax.set_ylabel('Amount (in millions)')
        ax.legend()
        plt.tight_layout()
        charts.append(fig)
    elif slide_name == "Competitor Analysis" and not custom_prompt:
        competitors = ['Competitor A', 'Competitor B', 'Our Company']
        market_share = [30, 45, 25]
        fig, ax = plt.subplots()
        ax.pie(market_share, labels=competitors, autopct='%1.1f%%')
        ax.set_title('Market Share Distribution')
        plt.tight_layout()
        charts.append(fig)
    else:
        # No predefined charts for custom slides
        pass
    return content, charts

def generate_data_driven_slide(uploaded_files, news_links, db_type, db_host, db_port, db_name, db_user, db_password, db_query):
    content = ""
    charts = []

    # Process uploaded files
    extracted_text = ""
    for uploaded_file in uploaded_files:
        if uploaded_file.type == "application/pdf":
            reader = PdfReader(uploaded_file)
            for page in reader.pages:
                extracted_text += page.extract_text()
        elif uploaded_file.type == "text/plain":
            extracted_text += uploaded_file.getvalue().decode("utf-8")

    # Process news links
    news_articles = ""
    links = news_links.strip().split('\n')
    for link in links:
        try:
            response = requests.get(link)
            soup = BeautifulSoup(response.content, 'html.parser')
            # Extract text from paragraphs
            paragraphs = soup.find_all('p')
            article_text = '\n'.join([para.get_text() for para in paragraphs])
            news_articles += article_text
        except Exception as e:
            st.write(f"Failed to fetch {link}: {e}")

    # Combine all text
    combined_text = extracted_text + "\n" + news_articles

    # Generate content summary using OpenAI

    if combined_text.strip():
        prompt = f"Summarize the following information for a pitch book presentation:\n{combined_text}"
        response = openai.chat.completions.create(
            model=MODEL,
            messages=[
                {"role": "system", "content": "You are an assistant helping with investment banking slides."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500,
            temperature=0.2,
        )
        content = response.choices[0].message.content.strip()        

        # Optionally, extract data points for charts
        # For simplicity, let's assume we have some data to plot
        data_points = [len(combined_text.split()), len(news_articles.split()), len(extracted_text.split())]
        labels = ['Total Words', 'News Articles', 'Uploaded Documents']
        fig, ax = plt.subplots()
        ax.bar(labels, data_points)
        ax.set_title('Data Overview')
        charts.append(fig)
    else:
        content = "No data provided."

    # Connect to SQL database and fetch data
    if db_query.strip():
        try:
            # Construct database URL
            if db_type == "MySQL":
                db_url = f"mysql+pymysql://{db_user}:{db_password}@{db_host}:{db_port}/{db_name}"
            elif db_type == "PostgreSQL":
                db_url = f"postgresql://{db_user}:{db_password}@{db_host}:{db_port}/{db_name}"
            elif db_type == "SQLite":
                db_url = f"sqlite:///{db_name}"
            else:
                db_url = ""

            if db_url:
                engine = create_engine(db_url)
                with engine.connect() as connection:
                    result = pd.read_sql_query(db_query, connection)
                    # Display fetched data
                    st.write("Data fetched from SQL database:")
                    st.dataframe(result)
                    # Generate chart from SQL data
                    if not result.empty:
                        fig_sql, ax_sql = plt.subplots()
                        result.plot(kind='bar', ax=ax_sql)
                        ax_sql.set_title('Database Data Visualization')
                        charts.append(fig_sql)
            else:
                st.write("Unsupported database type.")
        except Exception as e:
            st.write(f"Failed to connect to database: {e}")

    return content, charts

def chat_pod():
    st.header("Q&A bot")
    st.write("Ask any questions you have:")

    # Initialize session state for chat messages
    if 'chat_messages' not in st.session_state:
        st.session_state['chat_messages'] = []

    # Display previous messages
    for msg in st.session_state['chat_messages']:
        if msg['user'] == 'user':
            st.write(f"**You:** {msg['text']}")
        else:
            st.write(f"**Bot:** {msg['text']}")

    # User input
    user_input = st.text_input("Your message", key="chat_input")
    if st.button("Send", key="chat_send"):
        if user_input:
            # Append user message
            st.session_state['chat_messages'].append({'user': 'user', 'text': user_input})
            # Generate response
            response = openai.chat.completions.create(
                model=MODEL,
                messages=[
                    {"role": "system", "content": "You are a helpful assistant for investment banking queries."},
                    {"role": "user", "content": user_input}
                ],
                max_tokens=150,
                temperature=0.7,
            )
            bot_reply = response.choices[0].message.content.strip()
            st.session_state['chat_messages'].append({'user': 'bot', 'text': bot_reply})
            # Clear input
            #st.experimental_rerun()

def valuation_models():
    st.header("Valuation Models")
    st.write("Select the valuation models you want to use:")
    models = {
        "Discounted Cash Flow (DCF)": False,
        "Market Comparables": False,
        "Leveraged Buyout (LBO)": False,
        "Precedent Transactions": False
    }

    for model in models.keys():
        models[model] = st.checkbox(model, key=f"model_{model}")

    if st.button("Run Valuations", key="run_valuations"):
        selected_models = [model for model, selected in models.items() if selected]
        if selected_models:
            st.write("Running the following valuation models:")
            for model in selected_models:
                st.write(f"- {model}")
                # Run selected models
                if model == "Discounted Cash Flow (DCF)":
                    dcf_valuation()
                elif model == "Market Comparables":
                    market_comparables_valuation()
                elif model == "Leveraged Buyout (LBO)":
                    lbo_valuation()
                elif model == "Precedent Transactions":
                    precedent_transactions_valuation()
        else:
            st.write("Please select at least one valuation model.")

def dcf_valuation():
    st.subheader("Discounted Cash Flow (DCF) Valuation")
    st.write("Please input the following parameters:")
    revenue_growth = st.number_input("Revenue Growth Rate (%)", value=5.0, key="revenue_growth")
    ebit_margin = st.number_input("EBIT Margin (%)", value=15.0, key="ebit_margin")
    tax_rate = st.number_input("Tax Rate (%)", value=30.0, key="tax_rate")
    capex = st.number_input("Capital Expenditures (% of Revenue)", value=5.0, key="capex")
    nwc = st.number_input("Change in Net Working Capital (% of Revenue)", value=2.0, key="nwc")
    discount_rate = st.number_input("Discount Rate (%)", value=10.0, key="discount_rate")
    terminal_growth = st.number_input("Terminal Growth Rate (%)", value=2.0, key="terminal_growth")
    initial_revenue = st.number_input("Current Revenue (in millions)", value=100.0, key="initial_revenue")

    if st.button("Calculate DCF Valuation", key="calculate_dcf"):
        # Perform DCF calculations
        years = np.arange(1, 6)
        revenues = [initial_revenue * ((1 + revenue_growth/100) ** year) for year in years]
        EBIT = [rev * ebit_margin/100 for rev in revenues]
        NOPAT = [ebit * (1 - tax_rate/100) for ebit in EBIT]
        FreeCashFlow = [nopat - (rev * capex/100) - (rev * nwc/100) for nopat, rev in zip(NOPAT, revenues)]
        # Terminal Value
        TerminalValue = FreeCashFlow[-1] * (1 + terminal_growth/100) / ((discount_rate/100) - (terminal_growth/100))
        # Discount factors
        discount_factors = [(1 / ((1 + discount_rate/100) ** year)) for year in years]
        # Present Value of FCF
        PV_FCF = [fcf * df for fcf, df in zip(FreeCashFlow, discount_factors)]
        # PV of Terminal Value
        PV_TerminalValue = TerminalValue * discount_factors[-1]
        # Enterprise Value
        EV = sum(PV_FCF) + PV_TerminalValue
        st.write(f"The estimated Enterprise Value is: **${EV:.2f} million**")
        # Display results
        df = pd.DataFrame({
            'Year': years,
            'Revenue': revenues,
            'EBIT': EBIT,
            'NOPAT': NOPAT,
            'Free Cash Flow': FreeCashFlow,
            'Discount Factor': discount_factors,
            'PV of FCF': PV_FCF
        })
        st.dataframe(df.style.format({
            'Revenue': '${:,.2f}',
            'EBIT': '${:,.2f}',
            'NOPAT': '${:,.2f}',
            'Free Cash Flow': '${:,.2f}',
            'Discount Factor': '{:.4f}',
            'PV of FCF': '${:,.2f}'
        }))

# Placeholder functions for other valuation models
def market_comparables_valuation():
    st.subheader("Market Comparables Valuation")
    st.write("Functionality to be implemented.")

def lbo_valuation():
    st.subheader("Leveraged Buyout (LBO) Valuation")
    st.write("Functionality to be implemented.")

def precedent_transactions_valuation():
    st.subheader("Precedent Transactions Valuation")
    st.write("Functionality to be implemented.")

if __name__ == '__main__':
    main()
