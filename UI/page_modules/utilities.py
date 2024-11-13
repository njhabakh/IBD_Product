from pptx import Presentation
from pptx.util import Inches
import streamlit as st
import os
import datetime as dt

# Function to save content to a PowerPoint file
def save_content_to_ppt(filename="result/Profiler_Report.pptx"):
    # Ensure the result directory exists
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    
    # Create a PowerPoint presentation
    path = 'result/template_barclays.pptx'
    prs = Presentation(path)

    possible_slides = ["Overview", "Financials", "Geographic Mix", "Management Information",
                       "Recent News", "Discounted Cash Flow Analysis", "Leveraged Buyout Analysis"]

    # Create title slide
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    company_name_placeholder = slide.placeholders[0]
    date_placeholder = slide.placeholders[10]

    company_name_placeholder.text = st.session_state["company_ticker"] # Company name requires AI, use Ticker for now
    date_placeholder.text = dt.date.today().strftime("%Y-%m-%d")

    for title in possible_slides:
        if st.session_state.get(title):

            # Add a slide for each section
            slide = prs.slides.add_slide(prs.slide_layouts[15])
            title_placeholder = slide.placeholders[0]
            content_placeholder = slide.placeholders[14]

            # Set the title and content
            title_placeholder.text = title
            content_placeholder.text = st.session_state[title]

            # Add the chart to the Financials slide
            if title == "Financials" and os.path.exists(st.session_state.get("chart_path")):
                slide.shapes.add_picture(st.session_state["chart_path"], Inches(4), Inches(2), width=Inches(4))

    # Save the presentation
    prs.save(filename)
